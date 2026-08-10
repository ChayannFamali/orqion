"""Тесты парсинга документов (T-206).

Проверки:
- DOCX → markdown с заголовками (через Docling backend)
- MD/TXT → напрямую (без Docling)
- PDF fallback (pypdfium2) — текст без структуры заголовков, parser: "fallback"
- PDF без текстового слоя (скан) → failed с причиной
- parser помечается в результате: "docling" | "fallback" | "direct"
- DOCX с таблицей — таблица сохраняется в markdown
"""

from __future__ import annotations

from io import BytesIO
from typing import Any

import pytest
from app.rag.blob import BlobRef, BlobStore, LocalBlobStore
from app.rag.parser import parse_document


async def _put_blob(blob_store: BlobStore, content: bytes) -> BlobRef:
    """Кладёт байты в blob store и возвращает BlobRef."""

    async def _gen() -> Any:
        yield content

    return await blob_store.put(_gen())


def _make_docx() -> bytes:
    """Создаёт минимальный DOCX с заголовком и текстом."""
    from docx import Document as DocxDocument

    doc = DocxDocument()
    doc.add_heading("Test Heading", level=1)
    doc.add_paragraph("Hello orqion RAG")
    doc.add_heading("Subsection", level=2)
    doc.add_paragraph("Content here")
    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_text_pdf() -> bytes:
    """Создаёт PDF с текстовым слоем через reportlab."""
    try:
        from reportlab.pdfgen import canvas
    except ImportError:
        pytest.skip("reportlab not available for PDF test fixture")

    buf = BytesIO()
    c = canvas.Canvas(buf, pagesize=(595, 842))
    c.drawString(100, 700, "Hello from PDF text layer")
    c.drawString(100, 680, "Second line of text")
    c.save()
    return buf.getvalue()


def _make_scan_pdf() -> bytes:
    """Создаёт PDF без текстового слоя (эмуляция скана)."""
    try:
        from reportlab.pdfgen import canvas
    except ImportError:
        pytest.skip("reportlab not available for PDF test fixture")

    buf = BytesIO()
    c = canvas.Canvas(buf, pagesize=(595, 842))
    c.showPage()
    c.save()
    return buf.getvalue()


@pytest.fixture
def blob_store(tmp_path: str) -> LocalBlobStore:
    """Создаёт временный LocalBlobStore."""
    return LocalBlobStore(tmp_path)


@pytest.mark.asyncio
async def test_parse_docx(blob_store: LocalBlobStore) -> None:
    """DOCX → markdown с заголовками через Docling backend."""
    content = _make_docx()
    blob_ref = await _put_blob(blob_store, content)

    result = await parse_document(
        blob_store,
        sha256=blob_ref.sha256,
        filename="test.docx",
        blob_uri=blob_ref.uri,
    )

    assert result.error is None
    assert result.parser == "docling"
    assert "Test Heading" in result.markdown
    assert "Hello orqion RAG" in result.markdown
    assert "Subsection" in result.markdown


@pytest.mark.asyncio
async def test_parse_md_direct(blob_store: LocalBlobStore) -> None:
    """MD → напрямую, без Docling."""
    content = b"# Title\n\nSome markdown text\n"
    blob_ref = await _put_blob(blob_store, content)

    result = await parse_document(
        blob_store,
        sha256=blob_ref.sha256,
        filename="readme.md",
        blob_uri=blob_ref.uri,
    )

    assert result.error is None
    assert result.parser == "direct"
    assert "# Title" in result.markdown
    assert "Some markdown text" in result.markdown


@pytest.mark.asyncio
async def test_parse_txt_direct(blob_store: LocalBlobStore) -> None:
    """TXT → напрямую, без Docling."""
    content = b"Plain text content\nSecond line\n"
    blob_ref = await _put_blob(blob_store, content)

    result = await parse_document(
        blob_store,
        sha256=blob_ref.sha256,
        filename="notes.txt",
        blob_uri=blob_ref.uri,
    )

    assert result.error is None
    assert result.parser == "direct"
    assert "Plain text content" in result.markdown
    assert "Second line" in result.markdown


@pytest.mark.asyncio
async def test_parse_pdf_fallback(blob_store: LocalBlobStore) -> None:
    """PDF через fallback (pypdfium2) — текст без структуры, parser: "fallback"."""
    content = _make_text_pdf()
    blob_ref = await _put_blob(blob_store, content)

    result = await parse_document(
        blob_store,
        sha256=blob_ref.sha256,
        filename="doc.pdf",
        blob_uri=blob_ref.uri,
    )

    # Если ML-модели не установлены — fallback
    if result.parser == "fallback":
        assert result.error is None
        assert "Hello from PDF text layer" in result.markdown
    else:
        # ML-модели установлены — docling должен разобрать
        assert result.parser == "docling"
        assert result.error is None


@pytest.mark.asyncio
async def test_parse_pdf_scan_failed(blob_store: LocalBlobStore) -> None:
    """PDF без текстового слоя (скан) → failed с причиной."""
    content = _make_scan_pdf()
    blob_ref = await _put_blob(blob_store, content)

    result = await parse_document(
        blob_store,
        sha256=blob_ref.sha256,
        filename="scan.pdf",
        blob_uri=blob_ref.uri,
    )

    assert result.error is not None
    assert "текстового слоя" in result.error.lower() or "ocr" in result.error.lower()
    assert result.markdown == ""


@pytest.mark.asyncio
async def test_parse_docx_with_table(blob_store: LocalBlobStore) -> None:
    """DOCX с таблицей — таблица сохраняется в markdown."""
    from docx import Document as DocxDocument

    doc = DocxDocument()
    doc.add_heading("Report", level=1)
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "A"
    table.cell(1, 1).text = "42"
    buf = BytesIO()
    doc.save(buf)

    blob_ref = await _put_blob(blob_store, buf.getvalue())

    result = await parse_document(
        blob_store,
        sha256=blob_ref.sha256,
        filename="report.docx",
        blob_uri=blob_ref.uri,
    )

    assert result.error is None
    assert result.parser == "docling"
    assert "Report" in result.markdown
    assert "Name" in result.markdown
    assert "42" in result.markdown


@pytest.mark.asyncio
async def test_parse_pptx(blob_store: LocalBlobStore) -> None:
    """PPTX → markdown через Docling MsPowerpointDocumentBackend."""
    from pptx import Presentation as PptxPresentation

    prs = PptxPresentation()
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Slide Title"
    content = slide.placeholders[1]
    content.text = "Slide content text"

    buf = BytesIO()
    prs.save(buf)

    blob_ref = await _put_blob(blob_store, buf.getvalue())

    result = await parse_document(
        blob_store,
        sha256=blob_ref.sha256,
        filename="slides.pptx",
        blob_uri=blob_ref.uri,
    )

    assert result.error is None
    assert result.parser == "docling"
    assert "Slide Title" in result.markdown
    assert "Slide content text" in result.markdown


@pytest.mark.asyncio
async def test_parse_xlsx(blob_store: LocalBlobStore) -> None:
    """XLSX → markdown через Docling MsExcelDocumentBackend."""
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws["A1"] = "Name"
    ws["B1"] = "Value"
    ws["A2"] = "Alpha"
    ws["B2"] = "100"

    buf = BytesIO()
    wb.save(buf)

    blob_ref = await _put_blob(blob_store, buf.getvalue())

    result = await parse_document(
        blob_store,
        sha256=blob_ref.sha256,
        filename="data.xlsx",
        blob_uri=blob_ref.uri,
    )

    assert result.error is None
    assert result.parser == "docling"
    assert "Name" in result.markdown
    assert "Alpha" in result.markdown
    assert "100" in result.markdown
